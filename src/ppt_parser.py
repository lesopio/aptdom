"""
PPT解析器模块
负责从PPT文件中提取文本内容
"""

import logging
from pathlib import Path
from typing import Dict, List, Any, Union
from dataclasses import dataclass

from .logging_config import get_logger

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    get_logger(__name__).warning("python-pptx库未安装，PPT解析功能将受限")

logger = get_logger(__name__)


@dataclass
class SlideContent:
    """幻灯片内容数据类"""
    slide_index: int
    title: str
    text_content: str
    bullet_points: List[str]
    tables: List[Dict[str, Any]]
    images: List[Dict[str, Any]]
    notes: str


class PPTParser:
    """PPT解析器类"""
    
    def __init__(self):
        self.logger = get_logger(__name__)
    
    def extract_text(self, ppt_path: Union[str, Path]) -> List[SlideContent]:
        """
        从PPT文件中提取文本内容
        
        Args:
            ppt_path: PPT文件路径
            
        Returns:
            包含每页幻灯片内容的列表
        """
        ppt_path = Path(ppt_path)
        if not ppt_path.exists():
            raise FileNotFoundError(f"PPT文件不存在: {ppt_path}")
        
        self.logger.info(f"开始解析PPT文件: {ppt_path}")
        
        try:
            # 尝试使用python-pptx库解析
            return self._extract_with_pptx(ppt_path)
        except Exception as e:
            self.logger.warning(f"使用python-pptx解析失败: {e}")
            # 回退到基本文本提取
            return self._extract_basic_text(ppt_path)
    
    def _extract_with_pptx(self, ppt_path: Path) -> List[SlideContent]:
        """使用python-pptx库提取内容"""
        presentation = Presentation(ppt_path)
        slides_data = []
        
        for i, slide in enumerate(presentation.slides):
            self.logger.debug(f"处理第{i+1}页幻灯片")
            
            # 提取标题
            title = self._extract_slide_title(slide)
            
            # 提取文本内容
            text_content, bullet_points = self._extract_text_content(slide)
            
            # 提取表格
            tables = self._extract_tables(slide)
            
            # 提取图片信息
            images = self._extract_images(slide)
            
            # 提取备注
            notes = self._extract_notes(slide)
            
            slide_content = SlideContent(
                slide_index=i + 1,
                title=title,
                text_content=text_content,
                bullet_points=bullet_points,
                tables=tables,
                images=images,
                notes=notes
            )
            
            slides_data.append(slide_content)
        
        self.logger.info(f"成功提取{len(slides_data)}页幻灯片内容")
        return slides_data
    
    def _extract_basic_text(self, ppt_path: Path) -> List[SlideContent]:
        """基本文本提取（当python-pptx不可用时）"""
        self.logger.warning("使用基本文本提取模式，功能受限")
        
        # 这里可以添加其他PPT解析库的支持，如comtypes等
        # 目前返回一个占位符
        return [
            SlideContent(
                slide_index=1,
                title="PPT内容",
                text_content=f"PPT文件: {ppt_path.name}\n\n无法使用高级解析功能，请安装python-pptx库。",
                bullet_points=[],
                tables=[],
                images=[],
                notes=""
            )
        ]
    
    def _extract_slide_title(self, slide) -> str:
        """提取幻灯片标题"""
        if slide.shapes.title:
            return slide.shapes.title.text.strip()
        
        # 查找第一个文本框作为标题
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text.strip():
                return shape.text.strip()
        
        return f"幻灯片 {slide.slide_id}"
    
    def _extract_text_content(self, slide) -> tuple[str, List[str]]:
        """提取文本内容和项目符号"""
        text_content = ""
        bullet_points = []
        
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
                
            if shape == slide.shapes.title:
                continue  # 标题已单独处理
            
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                para_text = paragraph.text.strip()
                if not para_text:
                    continue
                
                # 检查是否为项目符号
                if paragraph.level > 0 or self._is_bullet_point(paragraph):
                    bullet_points.append(para_text)
                else:
                    text_content += para_text + "\n"
        
        return text_content.strip(), bullet_points
    
    def _is_bullet_point(self, paragraph) -> bool:
        """检查段落是否为项目符号"""
        # 检查段落是否有项目符号
        if paragraph.paragraph_format.bullet:
            return True
        
        # 检查文本是否以项目符号开头
        text = paragraph.text.strip()
        bullet_indicators = ["•", "-", "*", "→", "➢", "➤"]
        return any(text.startswith(indicator) for indicator in bullet_indicators)
    
    def _extract_tables(self, slide) -> List[Dict[str, Any]]:
        """提取表格内容"""
        tables = []
        
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table_data = []
                for row in shape.table.rows:
                    row_data = []
                    for cell in row.cells:
                        row_data.append(cell.text.strip())
                    table_data.append(row_data)
                
                tables.append({
                    "rows": len(table_data),
                    "cols": len(table_data[0]) if table_data else 0,
                    "data": table_data
                })
        
        return tables
    
    def _extract_images(self, slide) -> List[Dict[str, Any]]:
        """提取图片信息"""
        images = []
        
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                images.append({
                    "shape_id": shape.shape_id,
                    "name": shape.name,
                    "width": shape.width,
                    "height": shape.height,
                    "left": shape.left,
                    "top": shape.top
                })
        
        return images
    
    def _extract_notes(self, slide) -> str:
        """提取备注内容"""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            return slide.notes_slide.notes_text_frame.text.strip()
        return ""
