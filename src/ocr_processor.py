"""
OCR处理器模块
负责处理PPT中的图像内容
"""

import logging
import os
import tempfile
from pathlib import Path
from typing import Dict, List, Any, Optional
from dataclasses import dataclass

from .logging_config import get_logger

try:
    from PIL import Image
    import pytesseract
    from pptx import Presentation
except ImportError:
    get_logger(__name__).warning("OCR相关库未安装，OCR功能将受限")

logger = get_logger(__name__)


@dataclass
class OCRResult:
    """OCR结果数据类"""
    image_path: str
    text: str
    confidence: float
    bounding_box: tuple


class OCRProcessor:
    """OCR处理器类"""
    
    def __init__(self, tesseract_cmd: Optional[str] = None):
        self.logger = get_logger(__name__)
        
        # 设置tesseract路径（如果需要）
        if tesseract_cmd:
            pytesseract.pytesseract.tesseract_cmd = tesseract_cmd
        
        # 检查tesseract是否可用
        try:
            pytesseract.get_tesseract_version()
            self.tesseract_available = True
        except Exception as e:
            self.logger.warning(f"Tesseract OCR不可用: {e}")
            self.tesseract_available = False
    
    def process_slides(self, ppt_path: Path, slides_data: List[Any]) -> List[Any]:
        """
        处理PPT中的图像内容
        
        Args:
            ppt_path: PPT文件路径
            slides_data: 原始幻灯片数据
            
        Returns:
            包含OCR结果的更新后的幻灯片数据
        """
        if not self.tesseract_available:
            self.logger.warning("OCR功能不可用，跳过图像处理")
            return slides_data
        
        self.logger.info("开始OCR处理PPT中的图像...")
        
        # 创建临时目录用于存储提取的图像
        with tempfile.TemporaryDirectory() as temp_dir:
            temp_path = Path(temp_dir)
            
            # 提取PPT中的图像
            extracted_images = self._extract_images_from_ppt(ppt_path, temp_path)
            
            # 对每张图像执行OCR
            ocr_results = {}
            for image_info in extracted_images:
                slide_index = image_info["slide_index"]
                image_path = image_info["path"]
                
                try:
                    result = self._perform_ocr(image_path)
                    if result and result.text.strip():
                        if slide_index not in ocr_results:
                            ocr_results[slide_index] = []
                        ocr_results[slide_index].append(result)
                except Exception as e:
                    self.logger.error(f"OCR处理图像失败 {image_path}: {e}")
        
        # 更新幻灯片数据
        updated_slides = []
        for slide in slides_data:
            slide_index = slide.slide_index
            
            if slide_index in ocr_results:
                # 合并OCR结果到幻灯片
                updated_slide = self._merge_ocr_results(slide, ocr_results[slide_index])
                updated_slides.append(updated_slide)
            else:
                updated_slides.append(slide)
        
        self.logger.info(f"OCR处理完成，处理了{len(ocr_results)}页幻灯片")
        return updated_slides
    
    def _extract_images_from_ppt(self, ppt_path: Path, temp_dir: Path) -> List[Dict[str, Any]]:
        """从PPT中提取图像"""
        extracted_images = []
        
        try:
            presentation = Presentation(ppt_path)
            
            for i, slide in enumerate(presentation.slides):
                slide_index = i + 1
                image_count = 0
                
                for shape in slide.shapes:
                    if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                        try:
                            # 提取图像数据
                            image = shape.image
                            image_bytes = image.blob
                            
                            # 保存到临时文件
                            image_ext = image.ext or "png"
                            image_filename = f"slide_{slide_index}_img_{image_count}.{image_ext}"
                            image_path = temp_dir / image_filename
                            
                            with open(image_path, "wb") as f:
                                f.write(image_bytes)
                            
                            extracted_images.append({
                                "slide_index": slide_index,
                                "path": str(image_path),
                                "shape_id": shape.shape_id,
                                "name": shape.name,
                                "width": shape.width,
                                "height": shape.height
                            })
                            
                            image_count += 1
                            
                        except Exception as e:
                            self.logger.warning(f"提取图像失败 (slide {slide_index}): {e}")
        
        except Exception as e:
            self.logger.error(f"从PPT提取图像失败: {e}")
        
        self.logger.info(f"从PPT中提取了{len(extracted_images)}张图像")
        return extracted_images
    
    def _perform_ocr(self, image_path: str) -> Optional[OCRResult]:
        """对图像执行OCR"""
        try:
            # 打开图像
            image = Image.open(image_path)
            
            # 预处理图像以提高OCR准确性
            processed_image = self._preprocess_image(image)
            
            # 执行OCR
            ocr_data = pytesseract.image_to_data(
                processed_image, 
                output_type=pytesseract.Output.DICT
            )
            
            # 提取文本和置信度
            texts = []
            confidences = []
            bounding_boxes = []
            
            for i, text in enumerate(ocr_data["text"]):
                if text.strip():
                    texts.append(text)
                    confidences.append(float(ocr_data["conf"][i]))
                    bounding_boxes.append((
                        ocr_data["left"][i],
                        ocr_data["top"][i],
                        ocr_data["width"][i],
                        ocr_data["height"][i]
                    ))
            
            if texts:
                # 计算平均置信度
                avg_confidence = sum(confidences) / len(confidences)
                
                # 合并文本
                full_text = " ".join(texts)
                
                return OCRResult(
                    image_path=image_path,
                    text=full_text,
                    confidence=avg_confidence,
                    bounding_box=bounding_boxes[0] if bounding_boxes else (0, 0, 0, 0)
                )
            
        except Exception as e:
            self.logger.error(f"OCR处理失败 {image_path}: {e}")
        
        return None
    
    def _preprocess_image(self, image: Image.Image) -> Image.Image:
        """预处理图像以提高OCR准确性"""
        # 转换为灰度图
        if image.mode != "L":
            image = image.convert("L")
        
        # 调整大小（如果需要）
        width, height = image.size
        if width < 300 or height < 300:
            # 放大小图像
            scale = max(300 / width, 300 / height)
            new_size = (int(width * scale), int(height * scale))
            image = image.resize(new_size, Image.LANCZOS)
        
        # 增强对比度
        from PIL import ImageEnhance
        enhancer = ImageEnhance.Contrast(image)
        image = enhancer.enhance(1.5)
        
        return image
    
    def _merge_ocr_results(self, slide: Any, ocr_results: List[OCRResult]) -> Any:
        """将OCR结果合并到幻灯片数据中"""
        # 创建新的幻灯片数据
        updated_slide = type(slide)(
            slide_index=slide.slide_index,
            title=slide.title,
            text_content=slide.text_content,
            bullet_points=slide.bullet_points.copy(),
            tables=slide.tables.copy(),
            images=slide.images.copy(),
            notes=slide.notes
        )
        
        # 添加OCR结果到文本内容
        ocr_text = "\n\n[图像OCR识别结果]:\n"
        for i, result in enumerate(ocr_results, 1):
            ocr_text += f"图像{i}: {result.text}\n"
            ocr_text += f"置信度: {result.confidence:.2f}\n\n"
        
        # 更新文本内容
        if updated_slide.text_content:
            updated_slide.text_content += ocr_text
        else:
            updated_slide.text_content = ocr_text
        
        # 更新图像信息
        for result in ocr_results:
            # 查找对应的图像信息
            for img in updated_slide.images:
                if img.get("path", "").endswith(Path(result.image_path).name):
                    img["ocr_text"] = result.text
                    img["ocr_confidence"] = result.confidence
                    break
        
        return updated_slide
