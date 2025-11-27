"""
PPT解析器测试
"""

import unittest
from pathlib import Path
from unittest.mock import patch, MagicMock

from src.ppt_parser import PPTParser, SlideContent


class TestPPTParser(unittest.TestCase):
    """PPT解析器测试类"""
    
    def setUp(self):
        """测试前准备"""
        self.parser = PPTParser()
        
        # 创建测试PPT文件路径
        self.test_ppt = Path("tests/test_files/test_presentation.pptx")
        self.test_ppt.parent.mkdir(exist_ok=True)
        
        # 如果测试文件不存在，创建一个简单的测试文件
        if not self.test_ppt.exists():
            self._create_test_ppt()
    
    def _create_test_ppt(self):
        """创建测试PPT文件"""
        try:
            from pptx import Presentation
            from pptx.util import Inches
            
            # 创建演示文稿
            prs = Presentation()
            
            # 第1页：标题页
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title.text = "测试演示文稿"
            subtitle.text = "AI PPT转换器测试"
            
            # 第2页：内容页
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            title = slide.shapes.title
            body_shape = slide.shapes.placeholders[1]
            tf = body_shape.text_frame
            tf.text = "内容页"
            
            # 添加项目符号
            p = tf.add_paragraph()
            p.text = "第一点"
            p.level = 1
            
            p = tf.add_paragraph()
            p.text = "第二点"
            p.level = 1
            
            # 第3页：表格页
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # 添加标题
            left = top = Inches(1)
            width = Inches(2)
            height = Inches(0.5)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.text = "表格页"
            
            # 添加表格
            rows, cols = 3, 3
            left = Inches(1)
            top = Inches(2)
            width = Inches(6)
            height = Inches(2)
            table = slide.shapes.add_table(rows, cols, left, top, width, height).table
            
            # 填充表格
            for row in range(rows):
                for col in range(cols):
                    table.cell(row, col).text = f"单元格 {row+1},{col+1}"
            
            # 保存文件
            prs.save(str(self.test_ppt))
            
        except ImportError:
            # 如果python-pptx不可用，创建一个空文件
            self.test_ppt.write_text("测试PPT文件")
    
    @patch('src.ppt_parser.Presentation')
    def test_extract_text_with_pptx(self, mock_presentation):
        """测试使用python-pptx提取文本"""
        # 模拟Presentation对象
        mock_prs = MagicMock()
        mock_presentation.return_value = mock_prs
        
        # 模拟幻灯片
        mock_slide1 = MagicMock()
        mock_slide1.shapes.title.text = "标题1"
        mock_slide1.shapes = [
            MagicMock(has_text_frame=True, text="标题1"),
            MagicMock(has_text_frame=True, text="正文内容1"),
            MagicMock(has_text_frame=True, text="• 项目符号1"),
            MagicMock(has_text_frame=True, text="• 项目符号2"),
        ]
        mock_slide1.has_notes_slide = True
        mock_slide1.notes_slide.notes_text_frame.text = "备注内容1"
        
        mock_slide2 = MagicMock()
        mock_slide2.shapes.title.text = "标题2"
        mock_slide2.shapes = [
            MagicMock(has_text_frame=True, text="标题2"),
            MagicMock(has_text_frame=True, text="正文内容2"),
        ]
        mock_slide2.has_notes_slide = False
        
        mock_prs.slides = [mock_slide1, mock_slide2]
        
        # 执行测试
        result = self.parser.extract_text(self.test_ppt)
        
        # 验证结果
        self.assertEqual(len(result), 2)
        self.assertEqual(result[0].title, "标题1")
        self.assertEqual(result[0].text_content, "正文内容1")
        self.assertEqual(result[0].bullet_points, ["项目符号1", "项目符号2"])
        self.assertEqual(result[0].notes, "备注内容1")
        
        self.assertEqual(result[1].title, "标题2")
        self.assertEqual(result[1].text_content, "正文内容2")
        self.assertEqual(result[1].bullet_points, [])
        self.assertEqual(result[1].notes, "")
    
    def test_extract_basic_text(self):
        """测试基本文本提取"""
        # 模拟python-pptx导入失败
        with patch('src.ppt_parser.Presentation', side_effect=ImportError):
            result = self.parser.extract_text(self.test_ppt)
            
            # 验证结果
            self.assertEqual(len(result), 1)
            self.assertEqual(result[0].title, "PPT内容")
            self.assertIn("无法使用高级解析功能", result[0].text_content)
    
    def test_extract_slide_title(self):
        """测试提取幻灯片标题"""
        # 模拟幻灯片
        mock_slide = MagicMock()
        mock_slide.shapes.title = MagicMock()
        mock_slide.shapes.title.text = "测试标题"
        
        title = self.parser._extract_slide_title(mock_slide)
        self.assertEqual(title, "测试标题")
    
    def test_extract_text_content(self):
        """测试提取文本内容"""
        # 模拟幻灯片
        mock_slide = MagicMock()
        mock_slide.shapes.title = MagicMock()
        
        # 模拟形状
        mock_shape1 = MagicMock()
        mock_shape1.has_text_frame = True
        mock_shape1.text_frame.paragraphs = [
            MagicMock(text="段落1", level=0, paragraph_format=MagicMock(bullet=None)),
            MagicMock(text="• 项目符号1", level=1, paragraph_format=MagicMock(bullet=None)),
            MagicMock(text="段落2", level=0, paragraph_format=MagicMock(bullet=None)),
        ]
        
        mock_shape2 = MagicMock()
        mock_shape2.has_text_frame = True
        mock_shape2.text_frame.paragraphs = [
            MagicMock(text="• 项目符号2", level=1, paragraph_format=MagicMock(bullet=None)),
        ]
        
        mock_slide.shapes = [mock_slide.shapes.title, mock_shape1, mock_shape2]
        
        text_content, bullet_points = self.parser._extract_text_content(mock_slide)
        
        self.assertIn("段落1", text_content)
        self.assertIn("段落2", text_content)
        self.assertIn("项目符号1", bullet_points)
        self.assertIn("项目符号2", bullet_points)
    
    def test_is_bullet_point(self):
        """测试判断项目符号"""
        # 模拟段落
        mock_para = MagicMock()
        mock_para.paragraph_format.bullet = True
        mock_para.text = "测试项目符号"
        
        result = self.parser._is_bullet_point(mock_para)
        self.assertTrue(result)
        
        # 测试文本开头的项目符号
        mock_para.paragraph_format.bullet = None
        mock_para.text = "• 项目符号"
        result = self.parser._is_bullet_point(mock_para)
        self.assertTrue(result)
        
        # 测试普通文本
        mock_para.text = "普通文本"
        result = self.parser._is_bullet_point(mock_para)
        self.assertFalse(result)


if __name__ == "__main__":
    unittest.main()
