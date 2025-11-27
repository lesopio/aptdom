"""
简单使用示例
展示AI PPT转换器的基本用法
"""

import os
import sys
from pathlib import Path

# 添加项目路径
sys.path.insert(0, str(Path(__file__).parent.parent))

from src.main import main as ppt_converter_main


def create_sample_ppt():
    """创建一个示例PPT文件用于测试"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
        
        # 创建演示文稿
        prs = Presentation()
        
        # 第1页：标题页
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "AI PPT转换器演示"
        subtitle.text = "这是一个测试演示文稿"
        
        # 第2页：内容页
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        title = slide.shapes.title
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.text = "核心功能"
        
        # 添加项目符号
        p = tf.add_paragraph()
        p.text = "完整文本提取"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "AI智能处理"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "OCR图像识别"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "多种输出格式"
        p.level = 1
        
        # 第3页：表格页
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # 添加标题
        left = top = Inches(1)
        width = Inches(2)
        height = Inches(0.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = "功能对比"
        
        # 添加表格
        rows, cols = 4, 3
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(3)
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # 填充表格
        table.cell(0, 0).text = "功能"
        table.cell(0, 1).text = "传统转换"
        table.cell(0, 2).text = "AI转换"
        
        table.cell(1, 0).text = "文本提取"
        table.cell(1, 1).text = "✓"
        table.cell(1, 2).text = "✓"
        
        table.cell(2, 0).text = "内容整理"
        table.cell(2, 1).text = "✗"
        table.cell(2, 2).text = "✓"
        
        table.cell(3, 0).text = "OCR识别"
        table.cell(3, 1).text = "✗"
        table.cell(3, 2).text = "✓"
        
        # 第4页：总结页
        title_only_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(title_only_layout)
        title = slide.shapes.title
        title.text = "总结"
        
        # 添加文本框
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = "AI PPT转换器的优势："
        
        p = tf.add_paragraph()
        p.text = "• 智能内容整理"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "• 自动生成摘要"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "• 提取关键点"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "• 生成标签"
        p.level = 1
        
        # 保存文件
        sample_ppt = Path("sample_presentation.pptx")
        prs.save(str(sample_ppt))
        
        print(f"已创建示例PPT文件: {sample_ppt}")
        return sample_ppt
        
    except ImportError:
        print("python-pptx库未安装，无法创建示例PPT文件")
        print("请运行: pip install python-pptx")
        return None


def demo_basic_conversion():
    """演示基本转换功能"""
    print("\n=== 基本转换演示 ===")
    
    # 创建示例PPT
    sample_ppt = create_sample_ppt()
    if not sample_ppt:
        return
    
    # 保存原始sys.argv
    original_argv = sys.argv.copy()
    
    try:
        # 演示1: 基本转换 (Docx)
        print("\n1. 转换为Docx格式...")
        sys.argv = ['ppt-converter', str(sample_ppt), '--format', 'docx', '--verbose']
        ppt_converter_main()
        
        # 演示2: 转换为Markdown
        print("\n2. 转换为Markdown格式...")
        sys.argv = ['ppt-converter', str(sample_ppt), '--format', 'markdown', '--verbose']
        ppt_converter_main()
        
        # 演示3: 启用OCR
        print("\n3. 启用OCR处理...")
        sys.argv = ['ppt-converter', str(sample_ppt), '--ocr', '--verbose']
        ppt_converter_main()
        
        # 演示4: 使用OpenAI (如果配置了API密钥)
        api_key = os.getenv("API_KEY")
        if api_key:
            print("\n4. 使用OpenAI...")
            sys.argv = ['ppt-converter', str(sample_ppt), '--ai', 'openai', '--api-key', api_key, '--verbose']
            ppt_converter_main()
        else:
            print("\n4. 跳过OpenAI演示 (未设置API_KEY环境变量)")
        
        print("\n=== 演示完成 ===")
        print("生成的文件:")
        for ext in ['docx', 'md']:
            output_file = sample_ppt.with_suffix(f'.{ext}')
            if output_file.exists():
                print(f"  - {output_file}")
        
    except Exception as e:
        print(f"演示过程中发生错误: {e}")
    
    finally:
        # 恢复sys.argv
        sys.argv = original_argv


def demo_command_line():
    """演示命令行用法"""
    print("\n=== 命令行用法演示 ===")
    
    print("""
基本用法:
  ppt-converter presentation.pptx

指定输出格式:
  ppt-converter presentation.pptx --format markdown

指定输出文件:
  ppt-converter presentation.pptx -o output.docx

使用OpenAI:
  ppt-converter presentation.pptx --ai openai --api-key your_key

启用OCR:
  ppt-converter presentation.pptx --ocr

详细输出:
  ppt-converter presentation.pptx -v

组合使用:
  ppt-converter presentation.pptx --format markdown --ai openai --ocr -v

批量转换:
  python examples/batch_convert.py /path/to/ppt/files --format docx --workers 4
    """)


def main():
    """主函数"""
    print("AI PPT转换器 - 使用示例")
    print("=" * 50)
    
    # 演示基本转换
    demo_basic_conversion()
    
    # 演示命令行用法
    demo_command_line()
    
    print("\n提示:")
    print("1. 确保已安装所有依赖: pip install -r requirements.txt")
    print("2. 如果使用OCR，请安装Tesseract OCR")
    print("3. 如果使用OpenAI，请设置API密钥")
    print("4. 查看README.md获取更多详细信息")


if __name__ == "__main__":
    main()
